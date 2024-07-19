from typing import Optional, Iterable
import os
from pathlib import Path
import tempfile

from pptx import Presentation as NewPresentation
from pptx.util import Cm
from pptx.presentation import Presentation
from PIL.Image import Image

from core.data.caching import CardCache
from core.data.set_context import SetContext
from core.game_concepts.card import Card


class ImageSetPowerpoint:
    INCH_TO_CM = 2.54
    SCRYFALL_DPI = 96

    SLIDE_HEIGHT = 19.05
    SLIDE_WIDTH = 25.40
    MINIMUM_MARGIN = 2

    presentation: Optional[Presentation]

    @classmethod
    def from_image_list(cls, file_name: str, output_dir: str, images: Iterable[Image]):
        generator = cls(file_name, output_dir)
        for image in images:
            generator.add_centered_image_slide(image)
        return generator.create_powerpoint()

    def __init__(self, file_name: str, output_dir: str):
        self.presentation = NewPresentation()
        self.file_name = file_name
        self.output_dir = output_dir
        Path(self.output_dir).mkdir(parents=True, exist_ok=True)

    def create_powerpoint(self):
        file_path = os.path.join(self.output_dir, self.file_name)
        print(f"Saving to: {file_path}")
        self.presentation.save(file_path)
        return self.presentation

    def add_centered_image_slide(self, image: Image):
        image_width = (image.size[0] / self.SCRYFALL_DPI) * self.INCH_TO_CM
        image_height = (image.size[1] / self.SCRYFALL_DPI) * self.INCH_TO_CM
        width_ratio = image_width / (self.SLIDE_WIDTH - self.MINIMUM_MARGIN * 2)
        height_ratio = image_height / (self.SLIDE_HEIGHT - self.MINIMUM_MARGIN * 2)
        is_landscape = width_ratio >= height_ratio

        if is_landscape:
            new_image_width = image_width / width_ratio
            new_image_height = image_height / width_ratio
        else:
            new_image_width = image_width / height_ratio
            new_image_height = image_height / height_ratio

        x_position = (self.SLIDE_WIDTH - new_image_width) / 2
        y_position = (self.SLIDE_HEIGHT - new_image_height) / 2

        temporary_image_file = tempfile.NamedTemporaryFile(suffix=".png")
        image.save(temporary_image_file)

        blank_slide_layout = self.presentation.slide_layouts[6]
        slide = self.presentation.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(
            temporary_image_file, Cm(x_position), Cm(y_position), width=Cm(new_image_width), height=Cm(new_image_height)
        )


class PowerPointGenerator:
    set_context: SetContext

    @classmethod
    def create_set_review(cls, set_context: SetContext):
        generator = PowerPointGenerator(set_context)
        # TODO: Better handle output path logic.
        generator.generate_powerpoints(os.path.join(f'../../Generated Documents', set_context.set_code.upper()))
        return generator

    def __init__(self, set_context: SetContext):
        self.set_context = set_context

    def generate_powerpoints(self, output_dir: str = '.'):
        day_one_file_name = f"{self.set_context.set_code} - Commons and Uncommons.pptx"
        ImageSetPowerpoint.from_image_list(
            day_one_file_name,
            output_dir,
            (card.full_card_image for card in self.set_context.day_one_cards)
        )
        print(f"Created file '{day_one_file_name}'!")

        day_two_file_name = f"{self.set_context.set_code} - Rares and Mythics.pptx"
        ImageSetPowerpoint.from_image_list(
            day_two_file_name,
            output_dir,
            (card.full_card_image for card in self.set_context.day_two_cards)
        )
        print(f"Created file '{day_two_file_name}'!")

    @property
    def sorted_card_list(self) -> list[Card]:
        return self.set_context.sorted_card_list


def main(
        expansion: str,
        bonus_sheet: Optional[str],
        *queries: str,
        print_card_list: bool = False
) -> PowerPointGenerator:
    context = SetContext.from_queries(expansion, bonus_sheet, *queries, print_card_list=print_card_list)
    return PowerPointGenerator.create_set_review(context)


def otj():
    set_code = 'OTJ'
    bonus_set_code = 'OTP'
    scryfall_queries = [
        'set:otj unique:cards',
        'set:otp unique:cards',
        'set:big unique:cards',
        '(set:spg and date=otj) unique:cards'
    ]

    main(set_code, bonus_set_code, *scryfall_queries, print_card_list=True)


def mh3():
    set_code = 'MH3'
    bonus_set_code = None
    scryfall_queries = [
        "set:mh3 unique:cards",
        "(set:spg and date=mh3) unique:cards",
        "(game:paper) set:m3c t:legendary t:creature -is:reprint",
    ]

    main(set_code, bonus_set_code, *scryfall_queries, print_card_list=True)


def blb():
    set_code = 'BLB'
    bonus_set_code = None
    scryfall_queries = [
        "set:blb unique:cards cn<=261",
        "(set:spg and date=blb) unique:cards",
    ]

    main(set_code, bonus_set_code, *scryfall_queries, print_card_list=True)


def debug():
    file_name = "Test.pptx"
    cache = CardCache.from_card_keys(
        # TODO: Come up with more comprehensive test suite of cards.
        ("MOM", "Captive Weird"),
        ("MOM", "Invasion of Kaladesh"),
        ("MOM", "Joyful Stormsculptor"),
        ("HOU", "Refuse // Cooperate"),
        ("GRN", "Expansion // Explosion"),
        ("ELD", "Bonecrusher Giant"),
        ("NEO", "Fable of the Mirror Breaker"),
        ("ZNR", "Shatterskull Smashing"),
        ("CHK", "Akki Lavarunner"),
        ("BRO", "Skitterbeam Battalion"),
        ("EMN", "Hanweir Garrison"),  # TODO: Handle meld back faces
    )
    cards = cache.card_list()

    for card in cards:
        print(f"{card}: {card.layout}")

    output_dir = f'../../Generated Documents/Debug'
    path = Path(output_dir)
    path.mkdir(parents=True, exist_ok=True)

    ImageSetPowerpoint.from_image_list(
        file_name,
        output_dir,
        (card.full_card_image for card in cards)
    )
    print(f"Created file '{file_name}'!")


if __name__ == "__main__":
    blb()



