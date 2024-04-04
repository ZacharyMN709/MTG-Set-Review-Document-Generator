from typing import Iterable

from pptx import Presentation as NewPresentation
from pptx.util import Cm
from pptx.presentation import Presentation
from PIL.Image import Image
import tempfile

SLIDE_HEIGHT = 19.05
SLIDE_WIDTH = 25.40

SCRYFALL_DPI = 96
INCH_TO_CM = 2.54
MINIMUM_MARGIN = 2


def gen_powerpoint(
        file_name: str,
        card_images: Iterable[Image]
):
    prs = NewPresentation()
    for image in card_images:
        add_centered_image_slide(prs, image)
    prs.save(file_name)


def add_image_slide(
        presentation: Presentation,
        image_path: str,
        left_margin: float = 2,
        top_margin: float = 2,
        height: float = None,
        width: float = None
):
    blank_slide_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank_slide_layout)
    left = Cm(left_margin)
    top = Cm(top_margin)
    if height:
        height = Cm(height)
    if width:
        width = Cm(width)

    slide.shapes.add_picture(image_path, left, top, height=height, width=width)


def add_centered_image_slide(
        presentation: Presentation,
        image: Image
):
    image_width = (image.size[0] / SCRYFALL_DPI) * INCH_TO_CM
    image_height = (image.size[1] / SCRYFALL_DPI) * INCH_TO_CM
    width_ratio = image_width / (SLIDE_WIDTH - MINIMUM_MARGIN * 2)
    height_ratio = image_height / (SLIDE_HEIGHT - MINIMUM_MARGIN * 2)
    is_landscape = width_ratio >= height_ratio

    if is_landscape:
        new_image_width = image_width / width_ratio
        new_image_height = image_height / width_ratio
    else:
        new_image_width = image_width / height_ratio
        new_image_height = image_height / height_ratio

    x_position = (SLIDE_WIDTH - new_image_width) / 2
    y_position = (SLIDE_HEIGHT - new_image_height) / 2

    temporary_image_file = tempfile.NamedTemporaryFile(suffix=".png")
    image.save(temporary_image_file)

    blank_slide_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank_slide_layout)
    slide.shapes.add_picture(
        temporary_image_file, Cm(x_position), Cm(y_position), width=Cm(new_image_width), height=Cm(new_image_height)
    )
